"""Build the 2-hour NUS guest-lecture deck.

Usage:  python build_ppt.py
Output: nus_lecture_rfm_agent.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---- design tokens --------------------------------------------------------
NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2E, 0x86, 0xAB)
GREEN = RGBColor(0x06, 0xA7, 0x7D)
ORANGE = RGBColor(0xF1, 0x8F, 0x01)
GREY = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
BLANK = prs.slide_layouts[6]


# ---- helpers ---------------------------------------------------------------

def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def title_slide(title, subtitle, footer=""):
    s = prs.slides.add_slide(BLANK)
    _bg(s, NAVY)
    tf = _textbox(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.8))
    p = tf.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(40), True, WHITE
    tf2 = _textbox(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(1.0))
    p = tf2.paragraphs[0]
    p.text = subtitle
    p.font.size, p.font.color.rgb = Pt(20), RGBColor(0xBF, 0xD7, 0xEA)
    if footer:
        tf3 = _textbox(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.6))
        p = tf3.paragraphs[0]
        p.text = footer
        p.font.size, p.font.color.rgb = Pt(14), RGBColor(0x9F, 0xB7, 0xCA)
    return s


def section_slide(kicker, title, minutes=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, BLUE)
    tf = _textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.7))
    p = tf.paragraphs[0]
    p.text = kicker + (f"   ·   ~{minutes} min" if minutes else "")
    p.font.size, p.font.color.rgb = Pt(18), RGBColor(0xD9, 0xEC, 0xF4)
    tf2 = _textbox(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(2.0))
    p = tf2.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(38), True, WHITE
    return s


def content_slide(title):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    bar = s.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tf = _textbox(s, Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.9))
    p = tf.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(28), True, NAVY
    return s


def bullets(slide, items, left=0.8, top=1.6, width=11.7, height=5.4,
            size=18, color=None):
    """items: list of (level, text) or plain strings (level 0)."""
    tf = _textbox(slide, Inches(left), Inches(top), Inches(width), Inches(height))
    first = True
    for item in items:
        level, text = item if isinstance(item, tuple) else (0, item)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ("•  " if level == 0 else "–  ") + text
        p.level = level
        p.font.size = Pt(size if level == 0 else size - 2)
        p.font.color.rgb = color or (NAVY if level == 0 else GREY)
        p.space_after = Pt(10 if level == 0 else 5)
    return tf


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def table_slide(slide, headers, rows, left=0.7, top=1.7, width=12.0,
                row_h=0.5, font=14):
    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(left), Inches(top), Inches(width),
        Inches(row_h * (len(rows) + 1)))
    table = shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        para = cell.text_frame.paragraphs[0]
        para.font.size, para.font.bold, para.font.color.rgb = Pt(font), True, WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(font)
            para.font.color.rgb = NAVY
    return table


def quadrant_slide(slide):
    """2x2 segment grid used on the four-segments slide."""
    data = [
        ("High Value Customer", "R high · F high · M high",
         "Champions — retain & grow", BLUE, 0, 0),
        ("Important Customer", "R high · F low · M high",
         "Big baskets, no habit — develop", GREEN, 1, 0),
        ("Important Maintaining\nCustomer", "R low · F high · M high",
         "Loyal but lapsing — win back", ORANGE, 0, 1),
        ("Normal Customer", "everyone else",
         "Broad base — nurture cheaply", GREY, 1, 1),
    ]
    w, h = Inches(5.6), Inches(2.35)
    x0, y0 = Inches(0.9), Inches(1.8)
    for name, rule, play, color, col, row in data:
        left = x0 + col * (w + Inches(0.4))
        top = y0 + row * (h + Inches(0.35))
        shp = slide.shapes.add_shape(1, left, top, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.25)
        p = tf.paragraphs[0]
        p.text = name
        p.font.size, p.font.bold, p.font.color.rgb = Pt(20), True, WHITE
        p2 = tf.add_paragraph(); p2.text = rule
        p2.font.size, p2.font.color.rgb = Pt(15), WHITE
        p3 = tf.add_paragraph(); p3.text = play
        p3.font.size, p3.font.italic, p3.font.color.rgb = Pt(14), True, WHITE


def pipeline_slide(slide):
    """Horizontal 5-step agent pipeline diagram."""
    steps = [
        ("1. Ingest", "raw CSV /\ndatabase"),
        ("2. Clean", "dedupe, fix types,\ndrop refunds"),
        ("3. RFM Score", "quintiles 1–5\nper customer"),
        ("4. Segment", "4 business\ngroups"),
        ("5. Report", "strategy +\nKPIs per group"),
    ]
    w, h, gap = Inches(2.2), Inches(1.7), Inches(0.32)
    x, y = Inches(0.7), Inches(2.3)
    for i, (head, body) in enumerate(steps):
        left = x + i * (w + gap)
        shp = slide.shapes.add_shape(1, left, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = BLUE if i % 2 == 0 else NAVY
        shp.line.fill.background()
        tf = shp.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = head
        p.font.size, p.font.bold, p.font.color.rgb = Pt(17), True, WHITE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = body
        p2.font.size, p2.font.color.rgb = Pt(12), WHITE
        p2.alignment = PP_ALIGN.CENTER
        if i < len(steps) - 1:
            ar = slide.shapes.add_shape(13, left + w + Inches(0.02),
                                        y + Inches(0.62), gap - Inches(0.04),
                                        Inches(0.45))
            ar.fill.solid(); ar.fill.fore_color.rgb = ORANGE
            ar.line.fill.background()


# ===========================================================================
# SLIDES
# ===========================================================================

# 1 — title
s = title_slide(
    "AI Agents for Retail Analytics:\nRFM Customer Segmentation & Automated Reporting",
    "Guest Lecture · Data Analytics in Business · National University of Singapore",
    "TikTok Engineer (AI Agents) · PhD Student, NTU Singapore")
note(s, "Welcome everyone. Today: how a classic business-analytics technique "
        "(RFM) becomes a fully automated AI agent. Half concepts, half live demo.")

# 2 — about me
s = content_slide("About your lecturer")
bullets(s, [
    "Software engineer at TikTok — building AI agents for data analytics",
    (1, "Agents that clean data, run analyses, and write reports autonomously"),
    "PhD student at NTU Singapore",
    (1, "Research interests: LLM-based agents and applied data analytics"),
    "Why this lecture?",
    (1, "RFM is the single most-used segmentation method in retail CRM"),
    (1, "It is also a perfect first 'AI agent' project — you can build "
        "everything you see today yourselves"),
])
note(s, "Brief intro — emphasise the industry + research mix; invite questions "
        "any time.")

# 3 — agenda
s = content_slide("Today's roadmap (2 hours)")
table_slide(s, ["#", "Part", "Time"], [
    ["1", "The business problem: why segment customers?", "15 min"],
    ["2", "RFM fundamentals: scores and the four segments", "25 min"],
    ["3", "From analyst to AI agent: pipeline anatomy", "25 min"],
    ["—", "Break", "10 min"],
    ["4", "Live demo: the RFM agent end-to-end", "25 min"],
    ["5", "Marketing strategy per segment + KPIs", "10 min"],
    ["6", "Limitations, ethics & where this goes next · Q&A", "10 min"],
], row_h=0.62, font=16)

# ---- PART 1 ---------------------------------------------------------------
section_slide("Part 1", "The business problem:\nwhy segment customers?", 15)

# 5
s = content_slide("Retail's oldest marketing mistake")
bullets(s, [
    "“Blast” marketing: the same email, discount, and ad for every customer",
    (1, "A loyal VIP gets the same 10% coupon as someone who bought once in 2023"),
    "Three predictable failures:",
    (1, "Wasted spend — discounts given to people who would buy anyway"),
    (1, "Churn — your best customers feel like strangers"),
    (1, "Spam fatigue — low-intent customers unsubscribe"),
    "Marketing budgets are finite → you must decide WHO gets WHAT",
])
note(s, "Ask the class: who has unsubscribed from a brand's emails this month? "
        "Why?")

# 6
s = content_slide("Not all customers are equal")
bullets(s, [
    "A small share of customers usually drives most of the revenue",
    (1, "In our demo dataset: ~13% of customers generate ~40% of revenue"),
    (1, "Industry rule of thumb: the 'Pareto' 80/20 pattern"),
    "Acquiring a new customer costs 5–7× more than retaining one",
    "Implication: treat customers differently, deliberately",
    (1, "→ that requires a way to GROUP customers by behaviour"),
])

# 7
s = content_slide("Customer segmentation in one sentence")
bullets(s, [
    "Segmentation = splitting your customer base into groups that behave "
    "similarly, so each group can be served differently",
    "Common bases for segmentation:",
    (1, "Demographic — age, location, income"),
    (1, "Psychographic — lifestyle, values"),
    (1, "Behavioural — what people actually DO (purchases, visits, clicks)"),
    "Behavioural data wins in retail: it is objective, abundant, and "
    "already in your order database",
    "RFM is the classic behavioural segmentation — invented for catalogue "
    "mail-order in the 1990s, still everywhere today",
])

# ---- PART 2 ---------------------------------------------------------------
section_slide("Part 2", "RFM fundamentals", 25)

# 9
s = content_slide("R · F · M — three questions about every customer")
table_slide(s, ["Letter", "Question", "Metric", "Intuition"], [
    ["Recency", "How recently did they buy?",
     "Days since last purchase", "Recent buyers respond best"],
    ["Frequency", "How often do they buy?",
     "Number of orders", "Habit predicts future buying"],
    ["Monetary", "How much do they spend?",
     "Total spend ($)", "Past spend predicts future value"],
], row_h=0.85, font=16)
tf = bullets(s, ["All three come straight from the transaction log — "
                 "no surveys, no extra data collection needed"],
             top=5.0, size=16)

# 10
s = content_slide("Why does such a simple model survive?")
bullets(s, [
    "Interpretable — a marketing manager can explain every segment",
    "Cheap — three GROUP BY aggregations on data you already have",
    "Actionable — each segment maps directly to a campaign type",
    "Empirically strong — RFM variables remain top predictors of "
    "response rate in most retail models",
    "Limitations exist too (we'll return to them):",
    (1, "Ignores product category, margins, and acquisition channel"),
    (1, "Backward-looking — describes the past, doesn't forecast"),
])

# 11
s = content_slide("From transactions to RFM — a tiny example")
bullets(s, ["Transaction log (snapshot date: 1 June)"], top=1.45, size=15)
table_slide(s, ["customer", "order date", "amount"], [
    ["Alice", "28 May", "$120"], ["Alice", "10 May", "$80"],
    ["Alice", "2 Apr", "$95"], ["Ben", "30 May", "$400"],
    ["Cara", "20 Jan", "$150"], ["Cara", "5 Jan", "$210"],
], top=1.95, width=6.0, row_h=0.42, font=13)
bullets(s, ["After GROUP BY customer:"], left=7.1, top=1.45, width=5.6, size=15)
table_slide(s, ["customer", "R (days)", "F", "M"], [
    ["Alice", "4", "3", "$295"],
    ["Ben", "2", "1", "$400"],
    ["Cara", "132", "2", "$360"],
], left=7.1, top=1.95, width=5.6, row_h=0.5, font=13)
bullets(s, ["One row per customer — this is the table everything else "
            "builds on"], top=5.6, size=16)

# 12
s = content_slide("Scoring: from raw values to 1–5")
bullets(s, [
    "Raw values aren't comparable (days vs counts vs dollars) → convert "
    "each to a 1–5 score by quintile (5 equal-sized buckets)",
    (1, "Recency: LOWER days = BETTER = higher score"),
    (1, "Frequency & Monetary: HIGHER = BETTER = higher score"),
    "A customer becomes a 3-digit profile, e.g. R5-F4-M5",
    "Then call a dimension “high” if its score is above the population mean",
    "Why quintiles? Robust to outliers, always balanced, no magic thresholds",
])
note(s, "Quick check: a customer who bought yesterday, once ever, for $900 — "
        "what does their profile look like? (R5, F1, M5)")

# 13
s = content_slide("The four business segments")
quadrant_slide(s)
note(s, "This is the 'classic Chinese e-commerce CRM' simplification of the "
        "full 11-segment RFM grid — four groups are enough to act on.")

# 14
s = content_slide("Worked example — which segment?")
table_slide(s, ["Customer", "R score", "F score", "M score", "Segment?"], [
    ["#1: bought last week, 14 orders, $3.2k", "5", "5", "5", "?"],
    ["#2: bought yesterday, 2 orders, $900", "5", "1", "4", "?"],
    ["#3: last seen 7 months ago, 12 orders, $2.5k", "1", "5", "5", "?"],
    ["#4: last seen 5 months ago, 2 orders, $90", "2", "2", "1", "?"],
], row_h=0.7, font=15)
bullets(s, [
    "Take 2 minutes with your neighbour — classify all four",
    (1, "Answers: High Value · Important · Important Maintaining · Normal"),
], top=5.3, size=16)
note(s, "Reveal answers one by one; ask WHY for each. #3 is the interesting "
        "one — high historical value but going cold.")

# 15
s = content_slide("Quick check ✋")
bullets(s, [
    "Why do we score Recency in reverse (fewer days → higher score)?",
    "A customer's M score is 5 but F score is 1 — what kind of shopper "
    "is this? What's the risk of ignoring them?",
    "Why use quintiles instead of fixed cut-offs like “> $1,000 = high”?",
    "What happens to the segments if we re-run the analysis next quarter?",
], size=20)
note(s, "5 minutes of discussion. Q4 foreshadows the automation argument: "
        "segments drift, so the analysis must be repeated — perfect job for "
        "an agent.")

# ---- PART 3 ---------------------------------------------------------------
section_slide("Part 3", "From analyst to AI agent", 25)

# 17
s = content_slide("How this is done manually today")
bullets(s, [
    "A typical analyst workflow:",
    (1, "Export CSV → fix Excel formats → pivot tables → RFM formulas → "
        "PowerPoint for the marketing team"),
    (1, "Takes hours; repeated monthly; error-prone"),
    "Each step is judgement + mechanics:",
    (1, "Mechanics: dedupe rows, parse dates, group, score  → automatable"),
    (1, "Judgement: what counts as an outlier? which strategy fits? → "
        "rules or an LLM can encode much of it"),
    "If a workflow is repetitive, rule-describable, and data-driven — "
    "it is a candidate for an agent",
])

# 18
s = content_slide("What is an AI agent?")
bullets(s, [
    "A program that pursues a GOAL by deciding and executing a sequence "
    "of steps itself — not just answering one question",
    "Anatomy: Goal → Perceive (read data) → Decide (rules / LLM) → "
    "Act (clean, compute, write) → Report",
    "Spectrum of autonomy:",
    (1, "Script: fixed steps, no decisions"),
    (1, "Rule-based agent: fixed steps + encoded decisions  ← our cleaner"),
    (1, "LLM agent: language model chooses actions & writes text  ← our "
        "strategist (optional mode)"),
    "Key engineering idea: use the cheapest mechanism that does the job — "
    "LLMs only where language and judgement are needed",
])
note(s, "De-hype moment: most production 'agents' are 80% deterministic "
        "pipeline + 20% LLM at the judgement/wording layer.")

# 19
s = content_slide("Our agent: a 5-step pipeline")
pipeline_slide(s)
bullets(s, [
    "Deterministic where correctness matters (cleaning, scoring) — "
    "LLM where language matters (strategy narrative)",
    "Every step logs WHAT it did and WHY → auditability",
], top=4.7, size=16)

# 20
s = content_slide("Step 2 — automated data cleaning")
bullets(s, [
    "Real transaction data is dirty. Our sample dataset (on purpose):",
    (1, "duplicate rows · missing customer ids · 'not-a-date' values"),
    (1, "amounts like “$1,234.56” · refunds (negative) · $99,000 fat-fingers"),
    "The cleaning agent applies ordered rules and logs each action:",
    (1, "“Removed 55 exact duplicate rows”"),
    (1, "“Winsorised 27 extreme amounts above the 99.5th percentile”"),
    "Why log? Trust. A black-box cleaner is a liability in business analytics",
], size=17)

# 21
s = content_slide("Step 3 — the RFM engine (the code is short!)")
bullets(s, [
    "Three aggregations + two functions:",
    (1, "rfm = df.groupby('customer_id').agg(recency, frequency, monetary)"),
    (1, "score 1–5 by quintile:  pd.qcut(series, q=5)"),
    (1, "segment = compare R/F/M scores to their means"),
    "~60 lines of pandas — the analytics core is the EASY part",
    "The value of the agent is everything around it: validation, cleaning, "
    "logging, narrative, repeatability",
], size=18)

# 22
s = content_slide("Step 4 — strategy generation: rules vs LLM")
table_slide(s, ["", "Rule playbook", "LLM strategist (Claude)"], [
    ["How", "Hand-written strategy per segment", "Prompted with segment stats"],
    ["Output", "Same text every run", "Tailored, cites actual numbers"],
    ["Cost", "Free, instant, offline", "API call, ~seconds, needs key"],
    ["Risk", "Generic advice", "Must validate (hallucination)"],
    ["Use when", "Compliance-critical, stable", "Insight & wording matter"],
], row_h=0.62, font=14)
bullets(s, ["Our agent does BOTH: playbook always, LLM narrative on top "
            "when an API key is present (graceful degradation)"],
        top=5.4, size=16)

# 23
s = content_slide("The actual prompt we send to Claude")
tf = _textbox(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.6))
for i, line in enumerate([
    "You are a retail CRM strategist. Below is an RFM segment summary",
    "for a retail business (currency: SGD).",
    "",
    "<segment statistics table>",
    "",
    "For EACH segment, write a concise marketing strategy recommendation",
    "(3–5 sentences) that references the actual numbers (customer share,",
    "revenue share, recency, frequency, spend). Return ONLY a JSON object",
    "mapping the exact segment name to its recommendation string.",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size, p.font.name, p.font.color.rgb = Pt(15), "Courier New", NAVY
bullets(s, [
    "Prompt-engineering lessons: give a role, give the data, constrain "
    "the output format, demand grounding in the numbers",
], top=5.4, size=16)

# 24 — break
s = section_slide("☕ Break", "Back in 10 minutes", 10)

# ---- PART 4 ---------------------------------------------------------------
section_slide("Part 4", "Live demo:\nthe RFM agent end-to-end", 25)

# 26
s = content_slide("What to watch for in the demo")
bullets(s, [
    "1.  The dirty dataset — can you spot the problems the agent must fix?",
    "2.  The agent's step-by-step narration (cleaning log)",
    "3.  The customer map — four colours, four behaviours",
    "4.  The revenue-concentration chart (who pays the bills?)",
    "5.  The auto-generated strategy report (downloadable Markdown)",
    "Demo app: Streamlit UI — all code is shared with you after class",
], size=18)
note(s, "Run: streamlit run app.py — use sample dataset first; if internet "
        "allows, re-run with the Claude strategist toggle on.")

# 27
s = content_slide("The demo dataset (deliberately dirty)")
table_slide(s, ["customer_id", "order_id", "order_date", "amount"], [
    ["C00412", "ORD12345", "2026-05-28", "184.20"],
    ["  c00412  ", "ORD12345", "2026-05-28", "184.20"],
    ["C00077", "ORD13001", "not-a-date", "$1,204.50"],
    ["C00231", "ORD13208", "05 Mar 2026", "-95.00"],
    ["", "ORD13422", "2026-04-11", "98700.00"],
], row_h=0.55, font=14)
bullets(s, [
    "5,533 rows · ~800 customers · duplicates, refunds, missing ids, "
    "mixed formats, outliers — exactly what you meet in industry",
], top=5.2, size=16)

# 28
s = content_slide("Demo results — the headline numbers")
table_slide(s, ["Segment", "Customers", "Avg recency", "Avg spend",
                "Revenue share"], [
    ["High Value", "12.8%", "30 days", "$3,596", "40.0%"],
    ["Important", "9.1%", "32 days", "$835", "6.7%"],
    ["Important Maintaining", "17.9%", "208 days", "$2,715", "42.4%"],
    ["Normal", "60.2%", "185 days", "$208", "10.9%"],
], row_h=0.62, font=15)
bullets(s, [
    "⅓ of customers ≈ 90% of revenue — and the single biggest revenue "
    "block (42%) hasn't shopped in ~7 months: a burning platform",
], top=5.2, size=17)
note(s, "If the live demo fails, this slide carries the story.")

# ---- PART 5 ---------------------------------------------------------------
section_slide("Part 5", "Strategy per segment", 10)

# 30
s = content_slide("High Value Customer — retain & grow")
bullets(s, [
    "Diagnosis: recent, frequent, top spend — the engine of the business",
    "Goal: retention + advocacy (NOT discounts — they buy on value)",
    "Plays:",
    (1, "VIP loyalty tier, exclusive perks, early product access"),
    (1, "Personalised service: dedicated support, birthday gifts"),
    (1, "Referral rewards — they bring in customers like themselves"),
    "KPIs: repeat rate · share of wallet · NPS / referrals",
])

# 31
s = content_slide("Important Customer — build the habit")
bullets(s, [
    "Diagnosis: recent + big baskets, but only 1–2 orders → no habit yet",
    "The highest-UPSIDE segment: spend is proven, frequency is not",
    "Plays:",
    (1, "Cross-sell complements to the first purchase"),
    (1, "Time-boxed “come back in 30 days” incentives"),
    (1, "Points programmes where frequency earns the reward"),
    "KPIs: purchase frequency · days between orders · 60-day repeat rate",
])

# 32
s = content_slide("Important Maintaining — win back · Normal — nurture")
bullets(s, [
    "Important Maintaining (loyal but lapsing):",
    (1, "Win-back campaign with a meaningful incentive + churn survey"),
    (1, "Retarget their historical favourite categories"),
    (1, "High-value lapsed accounts → personal outreach, not email blasts"),
    (1, "KPIs: reactivation rate, win-back ROI"),
    "Normal (broad, low-value base):",
    (1, "Low-cost automated channels only; broad seasonal promos"),
    (1, "Watch for risers — promote them to higher-touch journeys"),
    (1, "KPIs: engagement, upgrade rate to better segments"),
], size=17)

# 33
s = content_slide("How would we know the strategies WORK?")
bullets(s, [
    "Segmentation is a hypothesis — campaigns are the experiment",
    "A/B test inside each segment:",
    (1, "Half of 'Important Maintaining' gets the win-back offer, half doesn't"),
    (1, "Compare reactivation rate → causal evidence, not vibes"),
    "Re-run the agent monthly and track segment MIGRATION:",
    (1, "Are Important customers becoming High Value? (good)"),
    (1, "Are High Value sliding into Maintaining? (alarm)"),
    "This is where automation pays off: the analysis is free to repeat",
])

# ---- PART 6 ---------------------------------------------------------------
section_slide("Part 6", "Limits, ethics & what's next", 10)

# 35
s = content_slide("Limitations & ethics — be a responsible analyst")
bullets(s, [
    "RFM limits:",
    (1, "Backward-looking; blind to margins, categories, seasonality"),
    (1, "New customers always look 'low value' — don't starve them"),
    "LLM-in-the-loop limits:",
    (1, "Hallucination: always validate generated numbers against the data"),
    (1, "Same input ≠ same output — log prompts & responses for audit"),
    "Ethics:",
    (1, "Price/offer discrimination across segments can cross legal lines"),
    (1, "PDPA (SG): purchase histories are personal data — minimise, secure"),
    (1, "'Normal customer' ≠ a person who deserves worse service"),
], size=17)

# 36
s = content_slide("Beyond this lecture")
bullets(s, [
    "Upgrades to the analytics:",
    (1, "RFM → CLV (customer lifetime value) prediction; BG/NBD models"),
    (1, "K-means clustering on R/F/M instead of fixed rules"),
    "Upgrades to the agent:",
    (1, "Schedule it monthly; let it email the report automatically"),
    (1, "Tool-using LLM agents: let the model decide which analysis to run"),
    "Try it yourself — the full project is yours:",
    (1, "agent/ (≈300 lines of pandas) + app.py (Streamlit) + sample data"),
    (1, "Swap in any transaction CSV — your e-commerce side project counts!"),
])

# 37
s = content_slide("Key takeaways")
bullets(s, [
    "1.  Customers are not equal — revenue concentrates; act on it",
    "2.  RFM turns a transaction log into 4 actionable segments with "
    "three GROUP BYs",
    "3.  An AI agent = deterministic pipeline for correctness + LLM for "
    "judgement & language",
    "4.  Always keep the audit trail: logs make agents trustworthy",
    "5.  Strategy without measurement is decoration — A/B test per segment",
], size=20)

# 38
s = title_slide("Thank you — questions?",
                "Slides, code and dataset: Documents/nus_lecture",
                "TikTok Engineer (AI Agents) · PhD Student, NTU Singapore")

out = Path(__file__).parent / "nus_lecture_rfm_agent.pptx"
prs.save(out)
print(f"Saved {len(prs.slides._sldIdLst)} slides -> {out}")
